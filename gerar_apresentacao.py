# -*- coding: utf-8 -*-
"""Gera a apresentacao 'Sindrome da Folha Azul' em .pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Paleta de cores (tons azulados, combinando com o tema)
AZUL_ESCURO = RGBColor(0x0D, 0x2B, 0x4E)
AZUL_MEDIO = RGBColor(0x1E, 0x5F, 0x8C)
AZUL_CLARO = RGBColor(0xA9, 0xD6, 0xE5)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_ESCURO = RGBColor(0x33, 0x33, 0x33)
VERMELHO = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
LARGURA = prs.slide_width
ALTURA = prs.slide_height

blank = prs.slide_layouts[6]


def fundo(slide, cor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = cor


def caixa(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_par(p, texto, tamanho, cor, negrito=False, align=PP_ALIGN.LEFT, italico=False):
    p.text = texto
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(tamanho)
    run.font.bold = negrito
    run.font.italic = italico
    run.font.color.rgb = cor
    run.font.name = "Calibri"
    return p


def faixa_titulo(slide, titulo, cor_faixa=AZUL_MEDIO):
    """Faixa superior com o titulo do slide."""
    shape = slide.shapes.add_shape(1, 0, 0, LARGURA, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor_faixa
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    set_par(p, titulo, 30, BRANCO, negrito=True, align=PP_ALIGN.CENTER)
    return shape


def nota_imagem(slide, texto):
    """Rodape com a sugestao de imagem."""
    tb, tf = caixa(slide, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.6))
    p = tf.paragraphs[0]
    set_par(p, "Imagem sugerida: " + texto, 13, AZUL_MEDIO, italico=True)


def bullets(tf, itens, tamanho=20, cor=CINZA_ESCURO, espaco=6):
    for i, item in enumerate(itens):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        set_par(p, item, tamanho, cor)
        p.space_after = Pt(espaco)


def conteudo_box(slide, top=Inches(1.6), height=Inches(4.8)):
    return caixa(slide, Inches(0.9), top, Inches(11.5), height)


# ------------------------------------------------------------------
# Slide 1 - Capa
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, AZUL_ESCURO)

_, tf = caixa(s, Inches(1), Inches(2.0), Inches(11.3), Inches(2.0))
p = tf.paragraphs[0]
set_par(p, "Sindrome da Folha Azul", 54, BRANCO, negrito=True, align=PP_ALIGN.CENTER)
p2 = tf.add_paragraph()
set_par(p2, "Uma doenca rara que afeta o sistema respiratorio", 24, AZUL_CLARO, align=PP_ALIGN.CENTER)

_, tf = caixa(s, Inches(1), Inches(4.6), Inches(11.3), Inches(2.0))
for i, txt in enumerate(["Aluno: (Seu nome)", "Turma: (Sua turma)", "Data: (Data da apresentacao)"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    set_par(p, txt, 20, BRANCO, align=PP_ALIGN.CENTER)
    p.space_after = Pt(6)

nota_imagem(s, "pessoa em ambiente de floresta com tons azulados.")

# ------------------------------------------------------------------
# Slide 2 - O que e
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, BRANCO)
faixa_titulo(s, "O que e a Sindrome da Folha Azul?")
_, tf = conteudo_box(s)
paras = [
    "A Sindrome da Folha Azul seria uma doenca causada pela exposicao prolongada a um fungo encontrado em florestas umidas. Segundo relatos ficticios, esse fungo liberaria particulas microscopicas capazes de afetar o organismo humano.",
    "A doenca receberia esse nome devido ao aparecimento de manchas azuladas na pele dos infectados.",
]
for i, txt in enumerate(paras):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    set_par(p, txt, 22, CINZA_ESCURO)
    p.space_after = Pt(14)
nota_imagem(s, "ilustracao de fungos em ambiente natural.")

# ------------------------------------------------------------------
# Slide 3 - Origem
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, BRANCO)
faixa_titulo(s, "Origem da doenca")
_, tf = conteudo_box(s)
paras = [
    "Acredita-se que a Sindrome da Folha Azul teria surgido em regioes de mata fechada e alta umidade.",
    "O fungo responsavel se desenvolveria em folhas em decomposicao e locais pouco iluminados.",
    "Essas condicoes favoreceriam a liberacao de esporos que poderiam ser inalados por pessoas que frequentassem essas areas.",
]
for i, txt in enumerate(paras):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    set_par(p, txt, 22, CINZA_ESCURO)
    p.space_after = Pt(14)
nota_imagem(s, "floresta umida.")

# ------------------------------------------------------------------
# Slide 4 - Sintomas
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, BRANCO)
faixa_titulo(s, "Principais sintomas")
_, tf = conteudo_box(s)
p = tf.paragraphs[0]
set_par(p, "Os sintomas comecariam de forma leve e aumentariam com o tempo.", 22, CINZA_ESCURO)
p.space_after = Pt(16)
sintomas = [
    "- Tosse persistente",
    "- Cansaco excessivo",
    "- Febre baixa",
    "- Dificuldade para respirar",
    "- Tontura ocasional",
    "- Manchas azuladas nos bracos e pernas",
]
for txt in sintomas:
    p = tf.add_paragraph()
    set_par(p, txt, 22, AZUL_MEDIO, negrito=True)
    p.space_after = Pt(8)
nota_imagem(s, "infografico medico.")

# ------------------------------------------------------------------
# Slide 5 - Transmissao
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, BRANCO)
faixa_titulo(s, "Como ocorreria a transmissao?")
_, tf = conteudo_box(s)
paras = [
    "A transmissao aconteceria pela inalacao dos esporos presentes no ambiente.",
    "Locais fechados, umidos e com pouca circulacao de ar apresentariam maior risco.",
    "A doenca nao seria transmitida de pessoa para pessoa.",
]
for i, txt in enumerate(paras):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    set_par(p, txt, 22, CINZA_ESCURO)
    p.space_after = Pt(14)
nota_imagem(s, "particulas microscopicas no ar.")

# ------------------------------------------------------------------
# Slide 6 - Diagnostico
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, BRANCO)
faixa_titulo(s, "Diagnostico")
_, tf = conteudo_box(s)
p = tf.paragraphs[0]
set_par(p, "O diagnostico seria realizado por meio de:", 22, CINZA_ESCURO)
p.space_after = Pt(14)
itens = [
    "- Exames de sangue",
    "- Avaliacao respiratoria",
    "- Analise das manchas na pele",
    "- Historico de exposicao a areas de mata umida",
]
for txt in itens:
    p = tf.add_paragraph()
    set_par(p, txt, 22, AZUL_MEDIO, negrito=True)
    p.space_after = Pt(8)
p = tf.add_paragraph()
set_par(p, "Os medicos buscariam identificar sinais da presenca do fungo no organismo.", 22, CINZA_ESCURO)
p.space_before = Pt(10)
nota_imagem(s, "medico analisando exames.")

# ------------------------------------------------------------------
# Slide 7 - Tratamento
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, BRANCO)
faixa_titulo(s, "Tratamento")
_, tf = conteudo_box(s)
p = tf.paragraphs[0]
set_par(p, "O tratamento incluiria:", 22, CINZA_ESCURO)
p.space_after = Pt(14)
itens = [
    "- Repouso",
    "- Hidratacao adequada",
    "- Alimentacao equilibrada",
    "- Medicamentos antifungicos especificos",
    "- Acompanhamento medico regular",
]
for txt in itens:
    p = tf.add_paragraph()
    set_par(p, txt, 22, AZUL_MEDIO, negrito=True)
    p.space_after = Pt(8)
p = tf.add_paragraph()
set_par(p, "A recuperacao dependeria da gravidade dos sintomas.", 22, CINZA_ESCURO)
p.space_before = Pt(10)
nota_imagem(s, "atendimento medico.")

# ------------------------------------------------------------------
# Slide 8 - Cura e prevencao
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, BRANCO)
faixa_titulo(s, "Cura e prevencao")
_, tf = conteudo_box(s)
p = tf.paragraphs[0]
set_par(p, "A Sindrome da Folha Azul teria cura quando tratada corretamente.", 22, CINZA_ESCURO)
p.space_after = Pt(14)
p = tf.add_paragraph()
set_par(p, "Formas de prevencao:", 22, CINZA_ESCURO, negrito=True)
p.space_after = Pt(10)
itens = [
    "- Uso de mascaras em areas de mata umida",
    "- Evitar locais com excesso de fungos",
    "- Manter ambientes ventilados",
    "- Procurar ajuda medica ao surgirem sintomas",
]
for txt in itens:
    p = tf.add_paragraph()
    set_par(p, txt, 22, AZUL_MEDIO, negrito=True)
    p.space_after = Pt(8)
nota_imagem(s, "pessoa usando mascara em area de floresta.")

# ------------------------------------------------------------------
# Slide 9 - A verdade (slide de impacto)
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, AZUL_ESCURO)
shape = s.shapes.add_shape(1, 0, 0, LARGURA, Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = VERMELHO
shape.line.fill.background()
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_par(tf.paragraphs[0], "A verdade sobre essa doenca", 30, BRANCO, negrito=True, align=PP_ALIGN.CENTER)

_, tf = caixa(s, Inches(1), Inches(2.0), Inches(11.3), Inches(4.2))
p = tf.paragraphs[0]
set_par(p, "Ate aqui, todas as informacoes apresentadas pareciam reais.", 24, AZUL_CLARO, align=PP_ALIGN.CENTER)
p.space_after = Pt(24)
p = tf.add_paragraph()
set_par(p, "Porem, a Sindrome da Folha Azul NAO EXISTE.", 32, BRANCO, negrito=True, align=PP_ALIGN.CENTER)
p.space_after = Pt(24)
p = tf.add_paragraph()
set_par(p, "Ela foi criada apenas para demonstrar como uma noticia falsa pode parecer verdadeira quando utiliza linguagem cientifica, sintomas detalhados e explicacoes convincentes.", 22, AZUL_CLARO, align=PP_ALIGN.CENTER)
nota_imagem(s, "simbolo de alerta ou fake news.")

# ------------------------------------------------------------------
# Slide 10 - Conclusao
# ------------------------------------------------------------------
s = prs.slides.add_slide(blank)
fundo(s, BRANCO)
faixa_titulo(s, "Conclusao")
_, tf = conteudo_box(s, height=Inches(3.2))
p = tf.paragraphs[0]
set_par(p, "As fake news podem enganar muitas pessoas, principalmente quando tratam de assuntos relacionados a saude.", 22, CINZA_ESCURO)
p.space_after = Pt(14)
p = tf.add_paragraph()
set_par(p, "Por isso, e importante:", 22, CINZA_ESCURO, negrito=True)
p.space_after = Pt(10)
itens = [
    "- Verificar as fontes da informacao.",
    "- Consultar profissionais da area.",
    "- Desconfiar de noticias sem comprovacao cientifica.",
    "- Buscar informacoes em sites confiaveis.",
]
for txt in itens:
    p = tf.add_paragraph()
    set_par(p, txt, 22, AZUL_MEDIO, negrito=True)
    p.space_after = Pt(8)

# Caixa de mensagem final
msg = s.shapes.add_shape(1, Inches(1.5), Inches(5.5), Inches(10.3), Inches(1.0))
msg.fill.solid()
msg.fill.fore_color.rgb = AZUL_MEDIO
msg.line.fill.background()
tf = msg.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_par(tf.paragraphs[0],
        '"Nem tudo o que parece cientifico e verdadeiro. Sempre confira a fonte da informacao."',
        20, BRANCO, negrito=True, italico=True, align=PP_ALIGN.CENTER)

caminho = r"c:\Users\USER\Desktop\Curso-TS\Sindrome_da_Folha_Azul.pptx"
prs.save(caminho)
print("Apresentacao salva em:", caminho)
print("Total de slides:", len(prs.slides._sldIdLst))
